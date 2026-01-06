import streamlit as st
import os
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import io

# --- Configuration ---
st.set_page_config(
    page_title="DocuConvert | Pro",
    page_icon="📊",
    layout="wide"
)

# --- Helper Function: Format File Size ---
def format_file_size(size_in_bytes):
    """Converts bytes to readable KB or MB"""
    if size_in_bytes < 1024:
        return f"{size_in_bytes} bytes"
    elif size_in_bytes < 1024 * 1024:
        return f"{size_in_bytes / 1024:.2f} KB"
    else:
        return f"{size_in_bytes / (1024 * 1024):.2f} MB"

# --- Conversion Logic (Pure Python) ---
def convert_pdf(file_stream):
    try:
        reader = PdfReader(file_stream)
        text = []
        for page in reader.pages:
            content = page.extract_text()
            if content:
                text.append(content)
        return "\n\n".join(text)
    except Exception as e:
        return f"Error reading PDF: {str(e)}"

def convert_docx(file_stream):
    try:
        doc = Document(file_stream)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        return f"Error reading DOCX: {str(e)}"

def convert_pptx(file_stream):
    try:
        prs = Presentation(file_stream)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        return f"Error reading PPTX: {str(e)}"

def convert_excel(file_stream):
    try:
        xls = pd.read_excel(file_stream, sheet_name=None)
        output = []
        for sheet_name, df in xls.items():
            output.append(f"### Sheet: {sheet_name}")
            output.append(df.to_markdown(index=False))
        return "\n\n".join(output)
    except Exception as e:
        return f"Error reading Excel: {str(e)}"

def convert_html(file_stream):
    try:
        soup = BeautifulSoup(file_stream, "html.parser")
        return soup.get_text(separator="\n")
    except Exception as e:
        return f"Error reading HTML: {str(e)}"

# --- Main Application ---
def main():
    st.title("📄 Universal Document to Markdown")
    st.markdown("Convert documents securely and analyze storage efficiency.")

    # [2] Upload Area
    uploaded_files = st.file_uploader(
        "Drag and drop files (PDF, DOCX, XLSX, PPTX, HTML)", 
        accept_multiple_files=True,
        type=['pdf', 'docx', 'xlsx', 'pptx', 'html', 'htm']
    )

    if uploaded_files:
        st.markdown("---")
        
        for uploaded_file in uploaded_files:
            file_name = uploaded_file.name
            file_ext = os.path.splitext(file_name)[1].lower()
            
            with st.expander(f"Processing: {file_name}", expanded=True):
                
                # 1. Calculate Original Size
                # seek(0,2) moves cursor to end to get size, then seek(0) resets for reading
                uploaded_file.seek(0, 2)
                original_size = uploaded_file.tell()
                uploaded_file.seek(0)
                
                # 2. Conversion
                converted_text = ""
                if file_ext == ".pdf":
                    converted_text = convert_pdf(uploaded_file)
                elif file_ext == ".docx":
                    converted_text = convert_docx(uploaded_file)
                elif file_ext == ".pptx":
                    converted_text = convert_pptx(uploaded_file)
                elif file_ext == ".xlsx":
                    converted_text = convert_excel(uploaded_file)
                elif file_ext in [".html", ".htm"]:
                    converted_text = convert_html(uploaded_file)
                else:
                    st.error(f"⚠️ Format {file_ext} not supported.")
                    continue

                if not converted_text or converted_text.strip() == "":
                    st.warning("⚠️ No text extracted.")
                else:
                    # 3. Calculate Converted Size (approx bytes in UTF-8)
                    converted_bytes = len(converted_text.encode('utf-8'))
                    
                    # 4. Create Tabs for View
                    tab1, tab2 = st.tabs(["📝 Preview & Download", "📊 File Size Comparison"])
                    
                    # --- Tab 1: Preview ---
                    with tab1:
                        st.text_area("Markdown Output", converted_text, height=250, key=f"text_{file_name}")
                        
                        base_name = os.path.splitext(file_name)[0]
                        c1, c2 = st.columns(2)
                        with c1:
                            st.download_button("📥 Download Markdown", converted_text, f"{base_name}.md")
                        with c2:
                            st.download_button("📄 Download Text", converted_text, f"{base_name}.txt")

                    # --- Tab 2: Comparison ---
                    with tab2:
                        # Calculate reduction percentage
                        if original_size > 0:
                            reduction = (1 - (converted_bytes / original_size)) * 100
                        else:
                            reduction = 0
                            
                        st.markdown("#### 📉 Storage Efficiency Analysis")
                        
                        # Create a clean data table
                        data = {
                            "Metric": ["Original File Size", "Converted Text Size", "Space Saved"],
                            "Value": [
                                format_file_size(original_size), 
                                format_file_size(converted_bytes),
                                f"{reduction:.1f}%"
                            ]
                        }
                        df_comparison = pd.DataFrame(data)
                        
                        # Display table
                        st.table(df_comparison)
                        
                        # Big visual callout
                        st.info(f"✨ The text version is **{reduction:.1f}% smaller** than the original file.")

if __name__ == "__main__":
    main()
